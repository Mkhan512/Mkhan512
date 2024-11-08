import pandas as pd
from pptx import Presentation
from pptx.util import Pt

# Load the Excel file and read the relevant sheets
excel_path = 'D:/path/to/your/excel/file/Commulative III  (28 Oct  To 3rd Nov), 2024 District Rawalpindi.xlsx'
ppt_template_path = 'D:/path/to/your/ppt/template/Day-6 ERM Presentation district (Distt Rawalpindi) -SNID-III Sep-24.pptx'


# Load PowerPoint template
ppt_template_path = '/mnt/data/Day-6 ERM Presentation district (Distt Rawalpindi) -SNID-III Sep-24.pptx'
prs = Presentation(ppt_template_path)

# Define function to replace placeholders in text frames
def replace_placeholder(slide, placeholder, text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if placeholder in shape.text_frame.text:
                shape.text_frame.text = text

# Process Daily Coverage Sheets (D-1 to D-6)
for day in range(1, 7):
    sheet_name = f'D-{day}'
    data = xls.parse(sheet_name)
    
    # Get the relevant slide (adjust slide indices as needed)
    slide = prs.slides[day - 1]  # Assuming each day's data goes to a separate slide

    # Insert daily coverage data
    tehsil = data['Unnamed: 0'].iloc[0]
    overall_target = data['Unnamed: 2'].sum()  # Example calculation
    coverage = data['Unnamed: 4'].sum()
    coverage_percentage = data['Unnamed: 5'].mean()

    # Replace placeholders on the slide
    replace_placeholder(slide, "Overall Target", str(overall_target))
    replace_placeholder(slide, "Coverage", str(coverage))
    replace_placeholder(slide, "%age", f"{coverage_percentage:.2f}%")

# Process Cumulative Data (Commulative NID Nov-23 sheet)
cumulative_data = xls.parse('Commulative NID Nov-23')
cumulative_slide = prs.slides[6]  # Assuming this is the 7th slide

# Extract cumulative metrics
cumulative_coverage = cumulative_data['Coverage'].sum()
cumulative_target = cumulative_data['Overall Target'].sum()
cumulative_percentage = (cumulative_coverage / cumulative_target) * 100

# Replace cumulative placeholders
replace_placeholder(cumulative_slide, "Cumulative Coverage", str(cumulative_coverage))
replace_placeholder(cumulative_slide, "Cumulative Target", str(cumulative_target))
replace_placeholder(cumulative_slide, "Cumulative %age", f"{cumulative_percentage:.2f}%")

# Process Vaccine Wastage (Vaccine Wastage RWP sheet)
wastage_data = xls.parse('Vaccine Wastage RWP')
wastage_slide = prs.slides[7]  # Assuming this is the 8th slide

# Example metrics from wastage data
total_given = wastage_data['Vaccine Given'].sum()
total_used = wastage_data['Used'].sum()
total_balance = wastage_data['Balance'].sum()
total_wastage = wastage_data['Wastage'].mean()

# Replace placeholders in the wastage slide
replace_placeholder(wastage_slide, "Total Given", str(total_given))
replace_placeholder(wastage_slide, "Total Used", str(total_used))
replace_placeholder(wastage_slide, "Total Balance", str(total_balance))
replace_placeholder(wastage_slide, "Wastage %", f"{total_wastage:.2f}%")

# Save the new PowerPoint
output_path = '/mnt/data/Updated_Presentation_Rawalpindi.pptx'
prs.save(output_path)
