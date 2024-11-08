import logging

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logging.debug("Debug message: Script started")

import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
from io import BytesIO

# Step 1: Load Excel Data
excel_file = excel_file = r'D:\NID-III.xlsx'  # Replace with your actual file path
data = pd.read_excel(excel_file, sheet_name=None)  # Load all sheets
print(data.head(10))
# Step 2: Load PowerPoint Template
ppt_template = 'D:\\Day-5.pptx'  # Replace with your file path
ppt = Presentation(ppt_template)

# Function to create charts
def create_chart(data, title):
    fig, ax = plt.subplots()
    data.plot(kind='bar', ax=ax)
    ax.set_title(title)
    buffer = BytesIO()
    plt.savefig(buffer, format='png')
    buffer.seek(0)
    plt.close(fig)
    return buffer

# Step 3: Populate Slides with Data
for slide_num, (sheet_name, df) in enumerate(data.items(), start=1):
    slide = ppt.slides[slide_num]  # Navigate to the specific slide
    # Customize the following steps to fit your slide layout
    # Example 1: Add Table Data
    table = slide.shapes.add_table(rows=len(df)+1, cols=len(df.columns), left=Inches(1), top=Inches(1), width=Inches(8), height=Inches(2.5)).table
    
    # Header Row
    for col_num, col_name in enumerate(df.columns):
        table.cell(0, col_num).text = str(col_name)
    
    # Data Rows
    for row_num, row in df.iterrows():
        for col_num, value in enumerate(row):
            table.cell(row_num + 1, col_num).text = str(value)
    
    # Example 2: Add a Chart to the Slide
    # Assuming we want to show coverage % or any key metric
    if 'Coverage' in df.columns:  # Check for a relevant column
        coverage_data = df[['UC', 'Coverage']]  # Adjust as per actual column names
        chart = create_chart(coverage_data.set_index('UC'), 'Coverage Chart')
        slide.shapes.add_picture(chart, Inches(1), Inches(4), width=Inches(6), height=Inches(3))
        
        chart.close()

# Step 4: Save the Presentation
ppt.save('automated_presentation_output.pptx')
import logging

logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')
logging.debug("Debug message: Script started")

